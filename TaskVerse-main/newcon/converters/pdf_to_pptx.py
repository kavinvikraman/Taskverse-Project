import fitz
from pptx import Presentation
from pptx.util import Inches, Pt

def convert_pdf_to_pptx(input_path, output_path):
    # Open PDF
    pdf_document = fitz.open(input_path)
    
    # Create PowerPoint presentation
    prs = Presentation()
    
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        text = page.get_text()
        
        # Add a slide for each page
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use title and content layout
        
        # Set title
        title = slide.shapes.title
        title.text = f"Page {page_num + 1}"
        
        # Add content
        content = slide.shapes.placeholders[1]
        tf = content.text_frame
        tf.text = text
        
        # Format text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(12)
    
    prs.save(output_path)
    pdf_document.close()
